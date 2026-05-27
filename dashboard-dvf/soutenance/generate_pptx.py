"""
Génère le support de soutenance (PowerPoint) — thème sombre assorti au dashboard.
Lancer :  python3.13 generate_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ICI = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(ICI, "img")

BG     = RGBColor(0x0E, 0x15, 0x25)
PANEL  = RGBColor(0x18, 0x22, 0x3A)
ACCENT = RGBColor(0xFF, 0x7A, 0x1A)
BLUE   = RGBColor(0x4F, 0x9C, 0xF9)
WHITE  = RGBColor(0xEC, 0xF1, 0xF8)
MUTED  = RGBColor(0x8B, 0x9B, 0xB4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp


def txt(slide, l, t, w, h, lines, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI", space=6):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        if isinstance(line, tuple):
            text, opts = line
        else:
            text, opts = line, {}
        run = p.add_run(); run.text = text
        f = run.font
        f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
        f.color.rgb = opts.get("color", color); f.name = font
    return tb


def accent_bar(slide):
    box(slide, 0, 0, Inches(0.18), H, fill=ACCENT)


def header(slide, kicker, title):
    accent_bar(slide)
    txt(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
        [(kicker.upper(), {"size": 13, "color": ACCENT, "bold": True})])
    txt(slide, Inches(0.6), Inches(0.8), Inches(12.2), Inches(0.9),
        [(title, {"size": 30, "color": WHITE, "bold": True})])
    box(slide, Inches(0.62), Inches(1.62), Inches(2.2), Pt(3), fill=ACCENT)


def footer(slide, n):
    txt(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        [("DVF·Insight — Dashboard immobilier 2025", {"size": 10, "color": MUTED})])
    txt(slide, Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.3),
        [(str(n), {"size": 10, "color": MUTED})], align=PP_ALIGN.RIGHT)


def bullets(slide, items, left=Inches(0.7), top=Inches(2.0), width=Inches(8.0), size=18, gap=14):
    tb = slide.shapes.add_textbox(left, top, width, Inches(4.6)); tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0; text = it
        if isinstance(it, tuple):
            text, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level = lvl
        run = p.add_run()
        run.text = ("▹  " if lvl == 0 else "•  ") + text
        run.font.size = Pt(size if lvl == 0 else size-3)
        run.font.color.rgb = WHITE if lvl == 0 else MUTED
        run.font.name = "Segoe UI"
    return tb


# ============================ SLIDE 1 — Titre ============================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, W, Inches(0.18), fill=ACCENT)
txt(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.5),
    [("🏠 DVF · INSIGHT", {"size": 18, "color": ACCENT, "bold": True})])
txt(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(2),
    [("Analyse du marché", {"size": 46, "color": WHITE, "bold": True}),
     ("immobilier français", {"size": 46, "color": ACCENT, "bold": True})], space=2)
txt(s, Inches(0.85), Inches(4.5), Inches(11), Inches(0.6),
    [("Dashboard interactif des Demandes de Valeurs Foncières — 1ᵉʳ semestre 2025",
      {"size": 20, "color": MUTED})])
box(s, Inches(0.85), Inches(5.25), Inches(3), Pt(3), fill=ACCENT)
txt(s, Inches(0.85), Inches(5.6), Inches(11), Inches(1),
    [("Hugo Berton", {"size": 22, "color": WHITE, "bold": True}),
     ("Yboost IA & Data — Bachelor 1 · Ynov · Soutenance juin 2026", {"size": 15, "color": MUTED})], space=4)

# ============================ SLIDE 2 — Sommaire ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "Sommaire", "Au programme")
cols = [
    ["1.  Contexte & problématique", "2.  Le jeu de données DVF", "3.  La démarche (pipeline)",
     "4.  Nettoyage des données"],
    ["5.  Feature engineering", "6.  Architecture technique", "7.  Démonstration",
     "8.  Insights clés"],
    ["9.  Focus Bouches-du-Rhône", "10. Choix & limites", "11. Conclusion", "12. Questions"],
]
for i, c in enumerate(cols):
    box(s, Inches(0.7+i*4.1), Inches(2.2), Inches(3.8), Inches(3.4), fill=PANEL)
    bullets(s, c, left=Inches(0.95+i*4.1), top=Inches(2.5), width=Inches(3.4), size=15, gap=16)
footer(s, 2)

# ============================ SLIDE 3 — Contexte ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "01 · Contexte", "Le défi : donner du sens à la donnée brute")
bullets(s, [
    "Les données DVF recensent TOUTES les transactions immobilières en France (open data).",
    "~400 000 ventes rien que sur le 1ᵉʳ semestre 2025 → impossible à lire « à l'œil nu ».",
    ("Problématique : comment rendre ces données compréhensibles en 10 secondes ?", 0),
    "Objectif : transformer un CSV brut en un tableau de bord clair qui raconte une histoire.",
], top=Inches(2.1), width=Inches(8.4))
box(s, Inches(9.3), Inches(2.1), Inches(3.4), Inches(3.3), fill=PANEL, line=ACCENT)
txt(s, Inches(9.5), Inches(2.4), Inches(3), Inches(3),
    [("EN CHIFFRES", {"size": 12, "color": ACCENT, "bold": True}),
     ("404 868", {"size": 34, "color": WHITE, "bold": True}),
     ("transactions brutes", {"size": 13, "color": MUTED}),
     ("10 colonnes", {"size": 20, "color": BLUE, "bold": True}),
     ("prix, surface, type, lieu…", {"size": 13, "color": MUTED})], space=8)
footer(s, 3)

# ============================ SLIDE 4 — Dataset ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "02 · Données", "Le jeu de données : dataset_dvf_s1_2025.csv")
bullets(s, [
    "Source : Demandes de Valeurs Foncières (DVF), publiées par data.gouv.fr.",
    "Colonnes clés : valeur_fonciere, surface_reelle, type_bien, commune, code_departement, date_mutation.",
    "Données réelles donc « sales » : valeurs manquantes, doublons, montants aberrants (jusqu'à 15 milliards €).",
    ("« Garbage In, Garbage Out » → un nettoyage rigoureux est indispensable.", 0),
], top=Inches(2.1), width=Inches(12))
footer(s, 4)

# ============================ SLIDE 5 — Pipeline ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "03 · Démarche", "Un pipeline ETL en 6 étapes (le fil rouge de l'année)")
steps = [("Collecte", "DVF + INSEE", BLUE), ("SQLite", "immobilier.db", ACCENT),
         ("Nettoyage", "Pandas", ACCENT), ("Fusion", "pd.merge", ACCENT),
         ("Features +\nAgrégation", "prix_m2,\nsegments, groupby", ACCENT),
         ("Dashboard", "Streamlit\n(filtres, carte)", BLUE)]
x = Inches(0.5)
for i, (t, d, col) in enumerate(steps):
    box(s, x, Inches(2.7), Inches(1.6), Inches(1.7), fill=PANEL, line=col)
    txt(s, x, Inches(2.82), Inches(1.6), Inches(0.8),
        [(t, {"size": 13, "color": WHITE, "bold": True})], align=PP_ALIGN.CENTER)
    txt(s, x, Inches(3.62), Inches(1.6), Inches(0.7),
        [(d, {"size": 10, "color": MUTED})], align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, x + Inches(1.52), Inches(3.0), Inches(0.45), Inches(0.8),
            [("→", {"size": 22, "color": ACCENT, "bold": True})], align=PP_ALIGN.CENTER)
    x = Emu(int(x) + int(Inches(2.0)))
txt(s, Inches(0.6), Inches(5.0), Inches(12.2), Inches(0.9),
    [("Chaque étape = une séance du module : collecte multi-sources (J1-2), base SQLite (J3),",
      {"size": 14, "color": MUTED}),
     ("nettoyage Pandas (J4), fusion DVF×INSEE (J5), feature engineering & agrégation (J6), dashboard (J7).",
      {"size": 14, "color": MUTED})], space=3)
footer(s, 5)

# ============================ SLIDE 6 — Nettoyage ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "04 · Nettoyage", "Du brut au propre (Pandas)")
bullets(s, [
    "Suppression des valeurs manquantes essentielles (valeur, surface) et des doublons.",
    "Surfaces nulles/aberrantes écartées (un prix au m² a besoin d'une surface valide).",
    "Montants aberrants filtrés : négatifs, nuls, ou délirants (> 50 M€).",
    "Typage de la date (datetime) ; département gardé en texte (« 01 », « 2A »).",
], top=Inches(2.1), width=Inches(8.4))
box(s, Inches(9.3), Inches(2.1), Inches(3.4), Inches(3.3), fill=PANEL, line=BLUE)
txt(s, Inches(9.5), Inches(2.4), Inches(3), Inches(3),
    [("RÉSULTAT", {"size": 12, "color": BLUE, "bold": True}),
     ("404 868", {"size": 26, "color": MUTED, "bold": True}),
     ("↓ nettoyage", {"size": 13, "color": MUTED}),
     ("373 038", {"size": 30, "color": WHITE, "bold": True}),
     ("lignes exploitables (92 %)", {"size": 13, "color": MUTED})], space=6)
footer(s, 6)

# ============================ SLIDE 7 — Feature engineering ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "05 · Feature Engineering", "Créer des variables qui ont du sens")
bullets(s, [
    "prix_m2 = valeur_fonciere / surface_reelle  → l'indicateur central de comparaison.",
    "annee & mois extraits de la date → analyse de l'évolution temporelle.",
    "segment_prix : Éco (< 4 000) · Standard (4 000–8 000) · Luxe (> 8 000 €/m²).",
    "Focus géographique : indicateur booléen « Bouches-du-Rhône » (département 13).",
    ("Choix clé : la médiane plutôt que la moyenne (résistante aux valeurs extrêmes).", 0),
], top=Inches(2.0), width=Inches(12))
footer(s, 7)

# ============================ SLIDE 8 — Architecture ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "06 · Technique", "Architecture : développé de A à Z")
box(s, Inches(0.7), Inches(2.1), Inches(5.7), Inches(3.4), fill=PANEL, line=ACCENT)
txt(s, Inches(0.95), Inches(2.35), Inches(5.2), Inches(3),
    [("BACK-END — Python", {"size": 17, "color": ACCENT, "bold": True}),
     ("• Collecte DVF + INSEE (multi-sources)", {"size": 14, "color": WHITE}),
     ("• Base SQLite (raw_dvf, raw_insee, clean)", {"size": 14, "color": WHITE}),
     ("• Pandas : nettoyage, fusion, features", {"size": 14, "color": WHITE}),
     ("• src/ modulaire + main.py orchestrateur", {"size": 14, "color": WHITE}),
     ("• Sortie : un JSON unique", {"size": 14, "color": WHITE})], space=8)
box(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(3.4), fill=PANEL, line=BLUE)
txt(s, Inches(7.15), Inches(2.35), Inches(5.2), Inches(3),
    [("DASHBOARD — Streamlit (interactif)", {"size": 17, "color": BLUE, "bold": True}),
     ("• Filtres live : département, type, segment,", {"size": 14, "color": WHITE}),
     ("  prix, mois + recherche de commune", {"size": 13, "color": MUTED}),
     ("• Carte, graphiques Altair, KPIs dynamiques", {"size": 14, "color": WHITE}),
     ("• Export CSV de la sélection", {"size": 14, "color": WHITE}),
     ("• (+ version HTML/CSS/JS fournie en bonus)", {"size": 13, "color": MUTED})], space=8)
footer(s, 8)

# ============================ SLIDE 9 — Démo vue d'ensemble ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "07 · Démonstration", "Le dashboard Streamlit (filtres en direct)")
img = os.path.join(IMG, "streamlit-apercu.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(1.4), Inches(1.95), width=Inches(10.5))
footer(s, 9)

# ============================ SLIDE 10 — Insights national ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "08 · Insights", "Lecture nationale du marché")
for i, (val, lab, col) in enumerate([
    ("2 667 €/m²", "prix médian national", ACCENT),
    ("Paris", "département le plus cher (9 447 €/m²)", BLUE),
    ("Creuse", "le plus abordable (1 474 €/m²)", BLUE),
]):
    box(s, Inches(0.7+i*4.05), Inches(2.1), Inches(3.75), Inches(1.5), fill=PANEL, line=col)
    txt(s, Inches(0.9+i*4.05), Inches(2.25), Inches(3.4), Inches(1.2),
        [(val, {"size": 26, "color": WHITE, "bold": True}), (lab, {"size": 12, "color": MUTED})], space=4)
bullets(s, [
    "La carte choroplèthe révèle d'un coup d'œil la fracture Paris / Côte d'Azur vs reste du territoire.",
    "Grâce à la fusion INSEE, on rapporte le volume à la population : les Hautes-Alpes ressortent comme le marché le plus dynamique (≈ 1 129 ventes / 100 000 hab.).",
], top=Inches(4.0), width=Inches(12))
footer(s, 10)

# ============================ SLIDE 11 — Insights segmentation ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "08 · Insights", "Segmentation & type de bien")
bullets(s, [
    "72 % des biens relèvent du segment « Éco » (< 4 000 €/m²) : le marché reste majoritairement abordable.",
    "Seuls ~7 % des biens dépassent 8 000 €/m² (« Luxe »), très concentrés sur Paris et la Côte d'Azur.",
    "Au m², l'appartement (3 462 €) se vend plus cher que la maison (2 226 €) au national…",
    "… mais l'inverse s'observe dans les Bouches-du-Rhône (effet maison + terrain en Provence).",
], top=Inches(2.0), width=Inches(12))
footer(s, 11)

# ============================ SLIDE 12 — Focus 13 ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "09 · Démonstration", "Carte interactive & focus régional")
img = os.path.join(IMG, "streamlit-map-final.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(2.4), Inches(1.95), width=Inches(8.5))
txt(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
    [("Couleur = prix moyen au m² · taille = volume. En filtrant sur le « 13 », on obtient "
      "instantanément le focus Bouches-du-Rhône.", {"size": 13, "color": MUTED})])
footer(s, 12)

# ============================ SLIDE 13 — Choix & limites ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "10 · Recul", "Choix d'analyse & limites")
box(s, Inches(0.7), Inches(2.1), Inches(5.7), Inches(3.5), fill=PANEL, line=ACCENT)
txt(s, Inches(0.95), Inches(2.35), Inches(5.2), Inches(3.1),
    [("MES CHOIX", {"size": 16, "color": ACCENT, "bold": True}),
     ("• Médiane vs moyenne sur les prix", {"size": 14, "color": WHITE}),
     ("• prix_m2 borné à 500–15 000 €/m²", {"size": 14, "color": WHITE}),
     ("• Classement limité à la métropole", {"size": 14, "color": WHITE}),
     ("• Seuils mini de transactions (fiabilité)", {"size": 14, "color": WHITE})], space=8)
box(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(3.5), fill=PANEL, line=BLUE)
txt(s, Inches(7.15), Inches(2.35), Inches(5.2), Inches(3.1),
    [("LIMITES & PISTES", {"size": 16, "color": BLUE, "bold": True}),
     ("• Données DVF parfois imparfaites (multi-lots)", {"size": 14, "color": WHITE}),
     ("• 1 seul semestre → pas de tendance annuelle", {"size": 14, "color": WHITE}),
     ("• Fusion possible avec données INSEE", {"size": 14, "color": WHITE}),
     ("• Filtres interactifs par département", {"size": 14, "color": WHITE})], space=8)
footer(s, 13)

# ============================ SLIDE 14 — Conclusion ============================
s = prs.slides.add_slide(BLANK); bg(s); header(s, "11 · Conclusion", "Ce que je retiens")
bullets(s, [
    "Un projet data complet, de la donnée brute au dashboard interactif.",
    "Maîtrise de la chaîne : collecte → nettoyage → feature engineering → agrégation → visualisation.",
    "Un livrable qui répond à la question de départ : comprendre le marché en 10 secondes.",
    ("Compétences clés : Pandas, structuration de projet, dataviz et storytelling.", 0),
], top=Inches(2.1), width=Inches(12))
footer(s, 14)

# ============================ SLIDE 15 — Questions ============================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, W, Inches(0.18), fill=ACCENT)
txt(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5),
    [("Merci de votre attention", {"size": 40, "color": WHITE, "bold": True}),
     ("Place à vos questions", {"size": 24, "color": ACCENT, "bold": True})], space=10)
txt(s, Inches(0.85), Inches(5.2), Inches(11), Inches(0.6),
    [("Hugo Berton — Dashboard immobilier DVF 2025 — Yboost IA & Data", {"size": 14, "color": MUTED})])

out = os.path.join(ICI, "soutenance_dashboard_dvf.pptx")
prs.save(out)
print("PPTX généré :", out, "—", len(prs.slides._sldIdLst), "slides")
